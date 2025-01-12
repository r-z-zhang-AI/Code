from pptx import Presentation
from pptx.util import Inches
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# 创建演示文稿
prs = Presentation()

# 添加标题幻灯片
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "公司年度报告"
slide1.placeholders[1].text = "2024 年"

# 添加内容幻灯片
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "年度数据分析"
slide2.placeholders[1].text = "1. 销售增长\n2. 用户增长\n3. 市场预测"

# 添加图片幻灯片
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
slide3.shapes.add_picture("/home/ruizhe/Github/mkdocs-site/docs/ESSAY/image-3.png", Inches(1), Inches(1), height=Inches(3))

# 添加表格幻灯片
slide4 = prs.slides.add_slide(prs.slide_layouts[5])
rows, cols = 3, 3
table = slide4.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(6), Inches(2)).table
for i in range(rows):
    for j in range(cols):
        table.cell(i, j).text = f"数据 {i+1}-{j+1}"

# 添加柱状图幻灯片
slide5 = prs.slides.add_slide(prs.slide_layouts[5])
chart_data = CategoryChartData()
chart_data.categories = ['一月', '二月', '三月']
chart_data.add_series('增长率', (2.5, 3.0, 4.2))
chart = slide5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(3), chart_data).chart

# 保存文件
prs.save("annual_report.pptx")
